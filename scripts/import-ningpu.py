#!/usr/bin/env python3
"""Build the Ningpu catalog from the official editable PowerPoint source."""

from __future__ import annotations

import argparse
import json
import os
import re
import sys
from datetime import datetime, timezone
from pathlib import Path

try:
    from pptx import Presentation
except ImportError as exc:  # pragma: no cover - environment guard
    raise SystemExit("python-pptx is required: python -m pip install python-pptx") from exc


REPO_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_SOURCE = Path(
    os.environ.get(
        "NINGPU_CATALOG_PATH",
        r"D:\cobioer-b2b\training-data\科佰生物\训练资料\科佰训练\江苏宁普天然病原体质控品目录PPT.pptx",
    )
)
DEFAULT_OUTPUT = REPO_ROOT / "src" / "data" / "ningpu-products.generated.json"
OFFICIAL_PAGE = "https://www.leadingmed.cn/product_xq/17.html"
SOURCE_TITLE = "江苏宁普天然病原体质控品目录（应用版，2025-09-06）"
PATHOGEN_TYPES_EN = {
    "G+细菌": "Gram-positive bacteria",
    "G-细菌": "Gram-negative bacteria",
    "抗酸阳性": "acid-fast bacteria",
    "真菌": "fungi",
    "DNA病毒": "DNA viruses",
    "RNA病毒": "RNA viruses",
    "MTBC": "MTBC",
    "NTM": "NTM",
    "MDR-TB": "MDR-TB",
}


PATHOGENS = {
    "甲型流感病毒H1N1": "Influenza A virus H1N1",
    "甲型流感病毒H3N2": "Influenza A virus H3N2",
    "甲型流感病毒H6N6": "Influenza A virus H6N6",
    "甲型流感病毒H7N8": "Influenza A virus H7N8",
    "甲型流感病毒H9N2": "Influenza A virus H9N2",
    "乙型流感病毒Victoria系": "Influenza B virus, Victoria lineage",
    "乙型流感病毒Yamagata系": "Influenza B virus, Yamagata lineage",
    "呼吸道合胞病毒A型": "Respiratory syncytial virus A",
    "呼吸道合胞病毒B型": "Respiratory syncytial virus B",
    "副流感病毒1型": "Human parainfluenza virus 1",
    "副流感病毒2型": "Human parainfluenza virus 2",
    "副流感病毒3型": "Human parainfluenza virus 3",
    "人偏肺病毒": "Human metapneumovirus",
    "冠状病毒": "Coronavirus",
    "鼻病毒1A型": "Rhinovirus 1A",
    "腺病毒01型": "Human adenovirus 01",
    "腺病毒02型": "Human adenovirus 02",
    "腺病毒03型": "Human adenovirus 03",
    "腺病毒05型": "Human adenovirus 05",
    "腺病毒07型": "Human adenovirus 07",
    "博卡病毒": "Human bocavirus",
    "肺炎链球菌": "Streptococcus pneumoniae",
    "嗜肺军团菌": "Legionella pneumophila",
    "百日咳杆菌": "Bordetella pertussis",
    "脑膜炎奈瑟菌": "Neisseria meningitidis",
    "卡他莫拉杆菌": "Moraxella catarrhalis",
    "肺炎克雷伯菌": "Klebsiella pneumoniae",
    "化脓性链球菌": "Streptococcus pyogenes",
    "流感嗜血杆菌": "Haemophilus influenzae",
    "结核分枝杆菌": "Mycobacterium tuberculosis",
    "铜绿假单胞菌": "Pseudomonas aeruginosa",
    "白喉棒杆菌": "Corynebacterium diphtheriae",
    "鲍曼不动杆菌": "Acinetobacter baumannii",
    "乙型链球菌": "Group B Streptococcus",
    "新生隐球菌": "Cryptococcus neoformans",
    "肺孢子菌": "Pneumocystis jirovecii",
    "烟曲霉": "Aspergillus fumigatus",
    "肺炎衣原体": "Chlamydia pneumoniae",
    "肺炎支原体": "Mycoplasma pneumoniae",
    "EB病毒": "Epstein-Barr virus",
    "人类免疫缺陷病毒1型": "Human immunodeficiency virus type 1",
    "登革病毒1型": "Dengue virus 1",
    "登革病毒2型": "Dengue virus 2",
    "登革病毒3型": "Dengue virus 3",
    "登革病毒4型": "Dengue virus 4",
    "流行性腮腺炎病毒": "Mumps virus",
    "日本乙型脑炎病毒": "Japanese encephalitis virus",
    "寨卡病毒": "Zika virus",
    "基孔肯雅病毒": "Chikungunya virus",
    "西尼罗病毒": "West Nile virus",
    "黄热病毒": "Yellow fever virus",
    "发热伴血小板减少综合征病毒": "Severe fever with thrombocytopenia syndrome virus",
    "乙型肝炎病毒": "Hepatitis B virus",
    "丙型肝炎病毒": "Hepatitis C virus",
    "柯萨奇病毒A6": "Coxsackievirus A6",
    "柯萨奇病毒A10": "Coxsackievirus A10",
    "柯萨奇病毒A16": "Coxsackievirus A16",
    "柯萨奇病毒B3": "Coxsackievirus B3",
    "肠病毒40型": "Enterovirus 40",
    "肠病毒41型": "Enterovirus 41",
    "肠病毒71型": "Enterovirus 71",
    "水痘带状疱疹病毒": "Varicella-zoster virus",
    "猴痘病毒": "Mpox virus",
    "痘苗病毒": "Vaccinia virus",
    "轮状病毒": "Rotavirus",
    "诺如病毒": "Norovirus",
    "金黄色葡萄球菌": "Staphylococcus aureus",
    "大肠杆菌": "Escherichia coli",
    "幽门螺旋杆菌": "Helicobacter pylori",
    "空肠弯曲菌": "Campylobacter jejuni",
    "粪肠球菌": "Enterococcus faecalis",
    "沙门氏菌": "Salmonella",
    "志贺菌": "Shigella",
    "巨细胞病毒": "Cytomegalovirus",
    "单纯疱疹病毒1型": "Herpes simplex virus type 1",
    "单纯疱疹病毒2型": "Herpes simplex virus type 2",
    "麻疹病毒": "Measles virus",
    "风疹病毒": "Rubella virus",
    "光滑念珠菌": "Candida glabrata",
    "白色念珠菌": "Candida albicans",
    "沙眼衣原体": "Chlamydia trachomatis",
    "淋病奈瑟氏菌": "Neisseria gonorrhoeae",
    "人型支原体": "Mycoplasma hominis",
    "解脲脲原体": "Ureaplasma urealyticum",
}

SPECIAL_NAMES = {
    "单耐药结核分枝杆菌RR-TB": "Rifampicin-resistant Mycobacterium tuberculosis (RR-TB)",
    "耐多药结核分枝杆菌MDR-TB": "Multidrug-resistant Mycobacterium tuberculosis (MDR-TB)",
    "广泛耐药肺炎克雷伯菌": "Extensively drug-resistant Klebsiella pneumoniae",
    "广泛耐药鲍曼不动杆菌": "Extensively drug-resistant Acinetobacter baumannii",
}


def clean(value: str) -> str:
    return re.sub(r"\s+", " ", value or "").strip()


def normalize_spec(value: str) -> str:
    value = clean(value).replace("×", "x").replace("X", "x")
    return re.sub(r"10e([0-9])", r"10^\1", value)


def normalize_code(value: str) -> str:
    return clean(value).translate(str.maketrans("⁰¹²³⁴⁵⁶⁷⁸⁹", "0123456789")).replace(" ", "")


def translate_name(name: str, product_type: str, code: str) -> str:
    name = clean(name)
    if product_type == "hpv-cell":
        match = re.search(r"HPV\s*([0-9]+[a-zA-Z]?)", name)
        genotype = match.group(1) if match else re.search(r"HPV([0-9]+[a-zA-Z]?)MO", code, re.I).group(1)
        return f"Human Papillomavirus (HPV {genotype}) Cellular Quality Control"
    if product_type == "multiplex":
        match = re.search(r"混合\s*([0-9]+)", name)
        return f"Respiratory Pathogen Multiplex Molecular Quality Control (Mix {match.group(1)})"
    if name in SPECIAL_NAMES:
        return f"{SPECIAL_NAMES[name]} Molecular Quality Control"

    suffixes = {
        "molecular": ("分子质控品", "Molecular Quality Control"),
        "antigen": ("抗原质控品", "Antigen Quality Control"),
    }
    suffix_zh, suffix_en = suffixes[product_type]
    stem = clean(name.replace(suffix_zh, ""))
    translated = PATHOGENS.get(stem)
    if not translated:
        raise ValueError(f"Missing English pathogen name for {name!r} ({code})")
    return f"{translated} {suffix_en}"


def level_from_code(code: str) -> tuple[str, str]:
    if code.endswith("-H"):
        return "高值", "High level"
    if code.endswith("-M"):
        return "中值", "Medium level"
    if code.endswith("-L"):
        return "低值", "Low level"
    if code.endswith("C1"):
        return "水平1", "Level 1"
    if code.endswith("C2"):
        return "水平2", "Level 2"
    return "", ""


def make_product(
    code: str,
    name: str,
    product_type: str,
    spec: str,
    coverage: str = "",
) -> dict:
    code, name, spec, coverage = normalize_code(code), clean(name), normalize_spec(spec), clean(coverage)
    level_zh, level_en = level_from_code(code)
    name_en = translate_name(name, product_type, code)
    type_labels = {
        "molecular": ("分子质控品", "Molecular QC", "分子检测质控", "Molecular testing QC"),
        "multiplex": ("多重分子质控品", "Multiplex molecular QC", "多重病原体检测", "Multiplex pathogen testing"),
        "antigen": ("抗原质控品", "Antigen QC", "抗原检测质控", "Antigen testing QC"),
        "hpv-cell": ("HPV细胞质控品", "HPV cellular QC", "HPV核酸检测", "HPV nucleic acid testing"),
    }
    type_zh, type_en, app_zh, app_en = type_labels[product_type]
    description = f"宁普医疗官方目录中的{name}，货号 {code}，{spec}。适用于专业体外诊断质量控制和方法学验证。"
    description_en = f"{name_en} from the official Ningpu catalog, catalog number {code}, supplied at {spec}. For professional IVD quality-control and assay-validation workflows."
    if coverage:
        description += f"覆盖范围：{coverage}。"
        description_en += " Contact our sales team for the full multiplex composition and configuration."
    tags = ["宁普医疗", type_zh]
    tags_en = ["Ningpu Diagnostics", type_en]
    if level_zh:
        tags.append(level_zh)
        tags_en.append(level_en)
    return {
        "id": code,
        "name": f"{name}（{level_zh}）" if level_zh and level_zh not in name else name,
        "nameEn": f"{name_en} - {level_en}" if level_en else name_en,
        "description": description,
        "descriptionEn": description_en,
        "tags": tags,
        "tagsEn": tags_en,
        "applications": [app_zh, "方法学验证", "IVD试剂开发"],
        "applicationsEn": [app_en, "Method validation", "IVD assay development"],
        "specs": spec,
        "specsEn": level_en if product_type == "antigen" and level_en else spec,
        "status": "Available by inquiry",
        "listPrice": "询价",
        "source": OFFICIAL_PAGE,
        "note": f"数据来源：{SOURCE_TITLE}。具体批次、包装与供应状态以销售确认为准。",
        "noteEn": "Source: Ningpu Natural Pathogen Quality Control Catalog, application edition, dated 2025-09-06. Confirm the current lot, pack size, and availability with sales.",
        "assayFormat": type_en,
        **(
            {
                "rawApplication": coverage,
                "rawApplicationEn": "Multiplex composition is listed in the official catalog; contact sales to confirm the required configuration.",
            }
            if coverage
            else {}
        ),
    }


def table_rows(table) -> list[list[str]]:
    return [[clean(cell.text) for cell in row.cells] for row in table.rows]


def slide_tables(presentation: Presentation, slide_number: int):
    return [shape.table for shape in presentation.slides[slide_number - 1].shapes if shape.has_table]


def parse_tiered_table(table, product_type: str) -> list[dict]:
    products, current_name, current_coverage = [], "", ""
    for cells in table_rows(table):
        if not any(cells) or len(cells) < 3:
            continue
        if cells[0]:
            current_name = cells[0]
        code = cells[1]
        if not code:
            continue
        is_mix = code.lower().startswith("mix")
        if is_mix and len(cells) >= 4 and cells[2]:
            current_coverage = cells[2]
        if is_mix:
            spec = " ".join(cells[3:])
            row_type = "multiplex"
        else:
            spec = " ".join(cells[2:])
            row_type = product_type
        products.append(make_product(code, current_name, row_type, spec, current_coverage if is_mix else ""))
    return products


def parse_molecular(presentation: Presentation) -> list[dict]:
    products = []
    for slide_number in range(8, 17):
        for table in slide_tables(presentation, slide_number):
            products.extend(parse_tiered_table(table, "molecular"))
    for table in slide_tables(presentation, 17):
        products.extend(parse_tiered_table(table, "molecular"))
    for table in slide_tables(presentation, 18):
        products.extend(parse_tiered_table(table, "molecular"))

    by_id: dict[str, dict] = {}
    for product in products:
        existing = by_id.get(product["id"])
        if not existing:
            by_id[product["id"]] = product
            continue
        if "01.0110-" not in product["id"]:
            raise ValueError(f"Unexpected duplicate Ningpu catalog number: {product['id']}")
        level_zh, level_en = level_from_code(product["id"])
        existing["name"] = f"耐药结核分枝杆菌分子质控品（RR-TB / MDR-TB，{level_zh}）"
        existing["nameEn"] = f"Drug-resistant Mycobacterium tuberculosis Molecular Quality Control (RR-TB / MDR-TB) - {level_en}"
        existing["note"] += " 官方目录对 RR-TB 与 MDR-TB 列示相同货号；具体耐药表型请在询价时确认。"
        existing["noteEn"] += " The official catalog lists the same number for RR-TB and MDR-TB; confirm the required resistance phenotype when requesting a quote."
        existing["description"] += " 官方目录对 RR-TB 与 MDR-TB 列示相同货号，具体耐药表型需销售确认。"
        existing["descriptionEn"] += " The official catalog lists the same number for RR-TB and MDR-TB; confirm the resistance phenotype when requesting a quote."
    return list(by_id.values())


def parse_ngs(presentation: Presentation) -> list[dict]:
    products = []
    english_names = {
        "NP-BALF-01": "Bronchoalveolar Lavage Fluid mNGS Quality Control",
        "NP-cfDNA-01": "Plasma cfDNA mNGS Quality Control",
        "NP-CSF-01": "Cerebrospinal Fluid mNGS Quality Control",
        "NP-TB-01": "Mycobacterium tuberculosis tNGS Quality Control",
    }
    for slide_number in (23, 24):
        for table in slide_tables(presentation, slide_number):
            rows = table_rows(table)
            code, name = rows[0][0], clean(rows[0][1])
            organisms = [row[2] for row in rows if len(row) > 3 and row[2] and row[2] not in ("人源细胞", "人源DNA")]
            types = sorted({row[3] for row in rows if len(row) > 3 and row[3] and "干扰" not in row[3]})
            types_en = [PATHOGEN_TYPES_EN.get(value, value) for value in types]
            name_en = english_names[code]
            assay = "tNGS" if code == "NP-TB-01" else "mNGS"
            coverage = "、".join(organisms)
            products.append(
                {
                    "id": code,
                    "name": name,
                    "nameEn": name_en,
                    "description": f"宁普医疗{name}，覆盖 {len(organisms)} 种目录内病原体/微生物：{coverage}。用于{assay}检测流程的质量控制和性能验证。",
                    "descriptionEn": f"Ningpu {name_en}, covering {len(organisms)} catalog-listed organisms across {', '.join(types_en)}. Designed for quality control and performance validation of {assay} workflows.",
                    "tags": ["宁普医疗", f"{assay}质控品", *types],
                    "tagsEn": ["Ningpu Diagnostics", f"{assay} QC", "Multiplex pathogen panel"],
                    "applications": [f"{assay}质控", "全流程验证", "IVD试剂开发"],
                    "applicationsEn": [f"{assay} QC", "End-to-end workflow validation", "IVD assay development"],
                    "specs": "500 µL/管；目录通用包装为 6 管/盒；干冰运输；-20°C 保存",
                    "specsEn": "500 uL/tube; 6 tubes/box; dry-ice shipping; store at -20 C",
                    "status": "Available by inquiry",
                    "listPrice": "询价",
                    "source": OFFICIAL_PAGE,
                    "note": f"数据来源：{SOURCE_TITLE}。浓度梯度、组合和批次以询价确认为准。",
                    "noteEn": "Source: Ningpu Natural Pathogen Quality Control Catalog, application edition, dated 2025-09-06. Confirm concentration levels, panel composition, and lot availability when requesting a quote.",
                    "assayFormat": assay,
                    "rawApplication": coverage,
                    "rawApplicationEn": f"Catalog panel covering {len(organisms)} organisms; contact sales for the current panel composition.",
                }
            )
    return products


def parse_catalog(source: Path) -> dict:
    presentation = Presentation(source)
    molecular = parse_molecular(presentation)
    antigen = []
    for slide_number in (20, 21):
        for table in slide_tables(presentation, slide_number):
            antigen.extend(parse_tiered_table(table, "antigen"))
    hpv = []
    for slide_number in (26, 27, 28):
        for table in slide_tables(presentation, slide_number):
            hpv.extend(parse_tiered_table(table, "hpv-cell"))
    ngs = parse_ngs(presentation)

    items = [
        {
            "id": "molecular-quality-controls",
            "name": "分子质控品",
            "nameEn": "Molecular Quality Controls",
            "description": "天然病原体分子质控品，覆盖单项、多重及耐药病原体检测",
            "descriptionEn": "Natural-pathogen molecular controls for singleplex, multiplex, and antimicrobial-resistance assays",
            "applications": ["分子检测质控", "方法学验证"],
            "applicationsEn": ["Molecular testing QC", "Method validation"],
            "products": molecular,
        },
        {
            "id": "antigen-quality-controls",
            "name": "抗原质控品",
            "nameEn": "Antigen Quality Controls",
            "description": "面向抗原检测流程的天然病原体质控品",
            "descriptionEn": "Natural-pathogen controls for antigen-testing workflows",
            "applications": ["抗原检测质控", "性能验证"],
            "applicationsEn": ["Antigen testing QC", "Performance validation"],
            "products": antigen,
        },
        {
            "id": "ngs-quality-controls",
            "name": "高通量测序质控品",
            "nameEn": "NGS Quality Controls",
            "description": "面向 mNGS 与 tNGS 全流程验证的多病原体质控组合",
            "descriptionEn": "Multipathogen QC panels for end-to-end mNGS and tNGS workflow validation",
            "applications": ["mNGS质控", "tNGS质控", "全流程验证"],
            "applicationsEn": ["mNGS QC", "tNGS QC", "End-to-end validation"],
            "products": ngs,
        },
        {
            "id": "hpv-cell-quality-controls",
            "name": "HPV细胞质控品",
            "nameEn": "HPV Cellular Quality Controls",
            "description": "覆盖多种 HPV 基因型及高、中、低浓度梯度的细胞质控品",
            "descriptionEn": "Cellular controls covering multiple HPV genotypes at high, medium, and low levels",
            "applications": ["HPV核酸检测", "全流程质控"],
            "applicationsEn": ["HPV nucleic acid testing", "End-to-end QC"],
            "products": hpv,
        },
    ]
    all_products = [product for item in items for product in item["products"]]
    ids = [product["id"] for product in all_products]
    if len(ids) != len(set(ids)):
        duplicates = sorted({code for code in ids if ids.count(code) > 1})
        raise ValueError(f"Duplicate product IDs after normalization: {duplicates}")
    return {
        "source": [OFFICIAL_PAGE],
        "sourceDocument": SOURCE_TITLE,
        "generatedAt": datetime.now(timezone.utc).isoformat(),
        "slideCount": len(presentation.slides),
        "productCount": len(all_products),
        "items": items,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--source", type=Path, default=DEFAULT_SOURCE)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--write", action="store_true", help="Write the generated JSON to the product data directory")
    args = parser.parse_args()
    if not args.source.exists():
        parser.error(f"Ningpu catalog source not found: {args.source}")
    catalog = parse_catalog(args.source)
    counts = {item["id"]: len(item["products"]) for item in catalog["items"]}
    print(json.dumps({"source": str(args.source), "productCount": catalog["productCount"], "groups": counts}, ensure_ascii=False, indent=2))
    if args.write:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(json.dumps(catalog, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        print(f"Wrote {args.output}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
